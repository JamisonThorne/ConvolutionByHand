from sklearn.metrics import confusion_matrix
import seaborn as sns
import numpy as np
import os
import matplotlib.pyplot as plt
from Blur import Blur_Kernel as blur
from Emboss import Emboss_Kernel as emboss
from Left_Sobel import Left_Sobel_Kernel as l_sobel
from Top_Sobel import Top_Sobel_Kernel as t_sobel
from Bottom_Sobel import Bottom_Sobel_Kernel as b_sobel
from Right_Sobel import Right_Sobel_Kernel as r_sobel
from Outline import Outline_Kernel as outline
from Identity import Identity_Kernel as identity
from Sharpen import Sharpen_Kernel as sharpen
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx import Presentation
from pptx.util import Inches, Pt
import io

def main():
    my_path = os.path.dirname(os.path.abspath(__file__))
    ppt_path = os.path.join(my_path,"ppts")
    clear_contents_of_directory(ppt_path)   # create directory if not created
    image_ds = get_my_image(my_path)   # pull image from specified path in function
    kernel_size = (3,3) # 3x3 kernel size for simplicity
    image_ds = rgb2gray(image_ds)   # convert rgb image to grayscale, no known built in function for this... Look more into later
    [rows, cols] = np.shape(image_ds)   # matrix size of image (512x512 in example)
    image_array = np.empty(kernel_size) # initialize empty kernel array for appending various kernels
    kernels = []
    kernels.append(['Blur',blur.kernel(np.empty(kernel_size))])      # blur
    kernels.append(['Emboss',emboss.kernel(np.empty(kernel_size))])    # emboss
    kernels.append(['Left Sobel',l_sobel.kernel(np.empty(kernel_size))])   # left sobel
    kernels.append(['Top Sobel',t_sobel.kernel(np.empty(kernel_size))])   # top sobel
    kernels.append(['Bottom Sobel',b_sobel.kernel(np.empty(kernel_size))])   # bot sobel
    kernels.append(['Right Sobel',r_sobel.kernel(np.empty(kernel_size))])   # right sobel
    kernels.append(['Outline',outline.kernel(np.empty(kernel_size))])   # outline
    kernels.append(['Identity',identity.kernel(np.empty(kernel_size))])  # identity
    kernels.append(['Sharpen',sharpen.kernel(np.empty(kernel_size))])   # sharpen
    heatmap_plot_ppt_generator(kernels,ppt_path,'kernels')
    convolved_ppt = Presentation()
    for k in range(0,len(kernels)): # loop through kernels
        convolved_image = []
        for i in range(0,rows): # loop through matrix size
            for j in range(0,cols): # loop through matrix size
                image_array = current_image_array(image_ds,kernel_size,i,j) #   Pull current image_array from input image. Note: not optimal and poorly written
                #better way to do above could be initializing an array with n+1xn+1 size and allocate an outer perimeter of zeros... Start matrix at 1:n,1:n.....then loop through... could work?
                convolution = [
                    image_array[0][0]*kernels[k][1][0][0],
                    image_array[1][0]*kernels[k][1][1][0],
                    image_array[2][0]*kernels[k][1][2][0],
                    image_array[0][1]*kernels[k][1][0][1],
                    image_array[1][1]*kernels[k][1][1][1],
                    image_array[2][1]*kernels[k][1][2][1],
                    image_array[0][2]*kernels[k][1][0][2],
                    image_array[1][2]*kernels[k][1][1][2],
                    image_array[2][2]*kernels[k][1][2][2]]    #   Multiply kernel by 3x3 array in image
                convolved_image.append(sum(convolution))    #   Sum the convolution for the new value in the convolved matrix
        final_image = np.reshape(convolved_image,(rows,cols))   # Reshape convolutions into a 512x512 matrix
        dual_plot_ppt_generator(image_ds,final_image,'Input',kernels[k][0],ppt_path,'convolution',convolved_ppt)
        
def dual_plot_ppt_generator(input1,input2,input1_string,input2_string,ppt_path,ppt_name,ppt_pres):
    assert np.shape(input1)==np.shape(input2),"Inputs must be the same size"
    fig,ax = plt.subplots()
    fig.suptitle('Convolution') # or plt.suptitle('Main title')
    current_figure = plt.gca()
    current_figure.axes.get_xaxis().set_visible(False)
    current_figure.axes.get_yaxis().set_visible(False)
    plt.subplot(1,2,1)
    plt.imshow(input1,cmap='gray')
    plt.title(input1_string,fontsize=20)
    plt.subplot(1,2,2)
    plt.imshow(input2,cmap='gray')
    plt.title(input2_string,fontsize=20)
    picture = io.BytesIO()
    ax.figure.savefig(picture,format='png')
    picture.seek(0)
    prs_slide = ppt_pres.slides.add_slide(ppt_pres.slide_layouts[6])
    prs_pic = prs_slide.shapes.add_picture(picture, Inches(-1), Inches(-0.65),width=Inches(12), height=Inches(9))
    ppt_pres.save(os.path.join(ppt_path,ppt_name+'.pptx'))
    plt.close()

def heatmap_plot_ppt_generator(inputs,ppt_path,ppt_name):   #may have to adjust vmax/vmin based on input data you are sending
    ppt_pres = Presentation()
    for i in range(0,len(inputs)):
        fig, ax = plt.subplots(figsize=(10,10))
        heatmap = sns.heatmap(
            inputs[i][1],
            cmap='Spectral',
            vmax=inputs[i][1].max(),
            vmin=inputs[i][1].min(),
            linewidths=0.25,
            linecolor='black',
            annot=True,
            annot_kws={"fontsize":20},
            fmt='g',
            cbar=False,
            square=True)
        ax.set_title(inputs[i][0],fontsize=20)
        picture = io.BytesIO()
        ax.figure.savefig(picture,format='png')
        picture.seek(0)
        prs_slide = ppt_pres.slides.add_slide(ppt_pres.slide_layouts[6])
        prs_pic = prs_slide.shapes.add_picture(picture, Inches(0.5), Inches(-0.75),width=Inches(9), height=Inches(9))
        ppt_pres.save(os.path.join(ppt_path,ppt_name+'.pptx'))
        plt.close()


def clear_contents_of_directory(current_directory):
    if not os.path.exists(current_directory):   # check if  directory does not exist
        os.makedirs(current_directory)  # create path to directory

def current_image_array(input,kernel_size,row,col):
    output = np.ones(kernel_size)
    [rows,cols] = np.shape(input)
    rows = rows-1
    cols = cols-1
    if row==0:
        output[0][0] = 0
        output[0][1] = 0
        output[0][2] = 0
    if col==0:
        output[0][0] = 0
        output[1][0] = 0
        output[2][0] = 0
    if row==rows:
        output[2][0] = 0
        output[2][1] = 0
        output[2][2] = 0
    if col==cols:
        output[0][2] = 0
        output[1][2] = 0
        output[2][2] = 0
    for i in range(0,kernel_size[0]):
        for j in range(0,kernel_size[1]):
            if not (output[i][j] == 0):
                try:
                    output[i][j] = input[row+(i-1)][col+(j-1)]
                except:
                    1
    return output

def rgb2gray(rgb):
    return np.dot(rgb[...,:3], [0.2989, 0.5870, 0.1140])

def get_my_image(my_path):
    my_path = os.path.join(my_path,"image")
    my_dir = os.listdir(my_path)
    my_image_path = os.path.join(my_path,my_dir[0])
    return plt.imread(my_image_path)

if __name__=='__main__':
    main()